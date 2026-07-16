"""Skill 4: SKU Calculator (PA Internal Only)

查询 PANW 产品 SKU 计算方式。
数据由管理员以 doc/slide 形式上传到 data/sku/ 目录。
Agent 根据用户提问从文档中提取相关内容回复。
仅对 PA 内部人员开放。
"""

import json
import hashlib
import re
import time
from pathlib import Path

DATA_DIR = Path(__file__).parent.parent / "data" / "sku"
DB_PATH = DATA_DIR / "sku_rules.json"
FILES_DIR = DATA_DIR / "files"

TOOL_DEFINITION = {
    "type": "function",
    "function": {
        "name": "query_sku",
        "description": (
            "Look up SKU calculation rules for PANW products. For internal staff only. "
            "Searches through uploaded SKU documents (docs, slides, PDFs) and structured data. "
            "The agent should interpret the document content to answer the user's question. "
            "Use this when an INTERNAL user asks about SKU numbers, licensing, pricing tiers, "
            "or how to calculate the right SKU for a customer."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "query": {
                    "type": "string",
                    "description": "Product or SKU question (e.g., 'PA-450 licensing', 'Prisma Access tier calculation', 'XSIAM pricing model')",
                },
                "customer_size": {
                    "type": "string",
                    "description": "Optional: customer size info (e.g., '5000 users', '3 sites')",
                },
            },
            "required": ["query"],
        },
    },
}


def _load_db() -> dict:
    DATA_DIR.mkdir(parents=True, exist_ok=True)
    if DB_PATH.exists():
        return json.loads(DB_PATH.read_text())
    return {"products": [], "rules": [], "notes": "", "files": []}


def _save_db(data: dict):
    DATA_DIR.mkdir(parents=True, exist_ok=True)
    DB_PATH.write_text(json.dumps(data, indent=2, ensure_ascii=False))


def add_product_sku(product: str, skus: list[dict], calculation_notes: str = "") -> dict:
    """Admin: add/update SKU info for a product (structured data)."""
    db = _load_db()
    existing = next((p for p in db["products"] if p["product"].lower() == product.lower()), None)
    if existing:
        existing["skus"] = skus
        existing["calculation_notes"] = calculation_notes
    else:
        db["products"].append({
            "product": product,
            "skus": skus,
            "calculation_notes": calculation_notes,
        })
    _save_db(db)
    return {"product": product, "sku_count": len(skus)}


def add_sku_file(filename: str, content: bytes, product: str, description: str = "") -> dict:
    """Admin: upload a doc/slide/PDF file containing SKU rules."""
    FILES_DIR.mkdir(parents=True, exist_ok=True)
    db = _load_db()

    # Upsert: remove existing with same original filename
    db["files"] = [f for f in db.get("files", []) if f.get("original_name") != filename]

    file_hash = hashlib.md5(content).hexdigest()[:8]
    safe_name = re.sub(r"[^\w.-]", "-", filename)[:80]
    stored_name = f"{file_hash}-{safe_name}"

    (FILES_DIR / stored_name).write_bytes(content)

    # Try to extract text content for search
    text_content = _extract_text(FILES_DIR / stored_name, filename)

    entry = {
        "original_name": filename,
        "stored_name": stored_name,
        "product": product,
        "description": description,
        "size_bytes": len(content),
        "uploaded_at": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime()),
        "text_content": text_content[:50000] if text_content else "",
    }
    db.setdefault("files", []).append(entry)
    _save_db(db)
    return entry


def _extract_text(file_path: Path, original_name: str) -> str:
    """Try to extract text from uploaded file."""
    ext = Path(original_name).suffix.lower()

    if ext in (".txt", ".md", ".csv"):
        try:
            return file_path.read_text(errors="ignore")
        except Exception:
            return ""

    if ext in (".pptx",):
        try:
            from pptx import Presentation
            prs = Presentation(str(file_path))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                texts.append(text)
            return "\n".join(texts)
        except Exception:
            return ""

    if ext in (".docx",):
        try:
            import zipfile
            from bs4 import BeautifulSoup
            with zipfile.ZipFile(str(file_path)) as z:
                with z.open("word/document.xml") as f:
                    soup = BeautifulSoup(f.read(), "html.parser")
                    return soup.get_text("\n", strip=True)
        except Exception:
            return ""

    return ""


def list_products() -> list[str]:
    """List all products with SKU data."""
    db = _load_db()
    return [p["product"] for p in db["products"]]


def list_files() -> list[dict]:
    """List all uploaded SKU files."""
    db = _load_db()
    return [{"original_name": f["original_name"], "product": f["product"],
             "description": f["description"], "stored_name": f["stored_name"],
             "uploaded_at": f.get("uploaded_at", "")}
            for f in db.get("files", [])]


def remove_file(stored_name: str) -> bool:
    """Remove an uploaded SKU file."""
    db = _load_db()
    files = db.get("files", [])
    new_files = [f for f in files if f["stored_name"] != stored_name]
    if len(new_files) == len(files):
        return False
    db["files"] = new_files
    _save_db(db)
    file_path = FILES_DIR / stored_name
    if file_path.exists():
        file_path.unlink()
    return True


def add_rule(rule: str, applies_to: str = "all"):
    """Admin: add a general SKU calculation rule."""
    db = _load_db()
    db["rules"].append({"rule": rule, "applies_to": applies_to})
    _save_db(db)


async def handle(arguments: dict) -> str:
    """Execute the SKU query skill."""
    query = arguments.get("query", "").strip().lower()
    customer_size = arguments.get("customer_size", "")

    if not query:
        return "Please provide a product name or SKU-related question."

    db = _load_db()

    if not db["products"] and not db["rules"] and not db.get("files"):
        return (
            "SKU database is empty. Admin needs to upload SKU documents or add structured data.\n"
            "Use the admin panel to upload doc/slide files with SKU rules."
        )

    lines = []

    # Search uploaded files (highest priority - richest content)
    file_matches = []
    for f in db.get("files", []):
        searchable = f"{f.get('product', '')} {f.get('description', '')} {f.get('original_name', '')} {f.get('text_content', '')}".lower()
        if query in searchable or any(kw in searchable for kw in query.split()):
            file_matches.append(f)

    if file_matches:
        lines.append("## SKU Documents Found\n")
        for f in file_matches[:3]:
            lines.append(f"### {f['original_name']} ({f['product']})")
            if f.get("description"):
                lines.append(f"Description: {f['description']}")
            text = f.get("text_content", "")
            if text:
                # Extract most relevant section
                relevant = _extract_relevant_section(text, query)
                lines.append(f"\n{relevant}\n")
            else:
                lines.append(f"\n(Binary file - content could not be extracted. File: {f['stored_name']})\n")
            lines.append("---")

    # Search structured products
    matches = []
    for prod in db["products"]:
        searchable = f"{prod['product']} {prod.get('calculation_notes', '')}".lower()
        sku_text = " ".join(s.get("sku", "") + " " + s.get("description", "") for s in prod.get("skus", []))
        searchable += " " + sku_text.lower()
        if query in searchable or any(kw in searchable for kw in query.split()):
            matches.append(prod)

    if matches:
        lines.append("\n## Structured SKU Data\n")
        for prod in matches[:5]:
            lines.append(f"### {prod['product']}\n")
            if prod.get("calculation_notes"):
                lines.append(f"{prod['calculation_notes']}\n")
            if prod.get("skus"):
                lines.append("| SKU | Description | Notes |")
                lines.append("|-----|-------------|-------|")
                for sku in prod["skus"]:
                    lines.append(f"| {sku.get('sku', '')} | {sku.get('description', '')} | {sku.get('notes', '')} |")
                lines.append("")

    # General rules
    applicable_rules = []
    for rule in db.get("rules", []):
        if rule["applies_to"] == "all" or query in rule["applies_to"].lower():
            applicable_rules.append(rule["rule"])

    if applicable_rules:
        lines.append("\n**General Rules:**")
        for rule in applicable_rules:
            lines.append(f"- {rule}")

    if customer_size:
        lines.append(f"\n*Customer context: {customer_size}*")
        lines.append("Please confirm the exact requirements with your SE team for accurate sizing.")

    if not lines:
        available = [p["product"] for p in db["products"]]
        file_products = list(set(f.get("product", "") for f in db.get("files", [])))
        all_products = list(set(available + file_products))
        return (
            f"No SKU information found for '{query}'.\n"
            f"Available products: {', '.join(all_products) if all_products else 'none'}\n"
            f"Try a product name like 'Prisma Access', 'PA-450', 'Cortex XDR'."
        )

    return "\n".join(lines)


def _extract_relevant_section(text: str, query: str, context_chars: int = 2000) -> str:
    """Extract the most relevant section from document text based on query."""
    lines = text.split("\n")
    query_keywords = query.lower().split()

    # Score each line by keyword matches
    scored_lines = []
    for i, line in enumerate(lines):
        line_lower = line.lower()
        score = sum(1 for kw in query_keywords if kw in line_lower)
        if score > 0:
            scored_lines.append((i, score))

    if not scored_lines:
        # Return beginning of document
        return text[:context_chars] + ("..." if len(text) > context_chars else "")

    # Get the highest scoring region
    scored_lines.sort(key=lambda x: x[1], reverse=True)
    best_line_idx = scored_lines[0][0]

    # Extract context around best match
    start = max(0, best_line_idx - 5)
    end = min(len(lines), best_line_idx + 20)
    section = "\n".join(lines[start:end])

    if len(section) > context_chars:
        section = section[:context_chars] + "..."

    return section
