"""Skill 5: TechDocs & Deployment Documentation

查询官方 TechDocs 和内部部署文档库。
支持上传 datasheet/文档文件和文本形式的内部文档。
"""

import json
import hashlib
import re
import time
from pathlib import Path

import httpx
from bs4 import BeautifulSoup

DATA_DIR = Path(__file__).parent.parent / "data" / "techdocs"
INTERNAL_DOCS_PATH = DATA_DIR / "internal_docs.json"
FILES_DIR = DATA_DIR / "files"

HEADERS = {
    "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36",
    "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
}

TECHDOCS_BASE = "https://docs.paloaltonetworks.com"

TOOL_DEFINITION = {
    "type": "function",
    "function": {
        "name": "query_techdocs",
        "description": (
            "Search PANW official TechDocs (docs.paloaltonetworks.com) and internal deployment documentation. "
            "Also searches uploaded technical documents (PDFs, docs, slides). "
            "Use this when a user asks about product deployment, configuration, troubleshooting, "
            "best practices, or any technical documentation."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "query": {
                    "type": "string",
                    "description": "Deployment or technical question (e.g., 'PAN-OS 11.1 upgrade steps', 'Prisma Access GlobalProtect setup')",
                },
                "product": {
                    "type": "string",
                    "enum": ["panos", "panorama", "prisma-access", "prisma-cloud", "cortex-xdr", "cortex-xsiam", "cortex-xsoar", "cn-series", "vm-series", "all"],
                    "description": "Filter by product area",
                },
            },
            "required": ["query"],
        },
    },
}


def _load_internal_docs() -> list[dict]:
    DATA_DIR.mkdir(parents=True, exist_ok=True)
    if INTERNAL_DOCS_PATH.exists():
        return json.loads(INTERNAL_DOCS_PATH.read_text()).get("documents", [])
    return []


def _save_internal_docs(docs: list[dict]):
    DATA_DIR.mkdir(parents=True, exist_ok=True)
    INTERNAL_DOCS_PATH.write_text(json.dumps({"documents": docs}, indent=2, ensure_ascii=False))


def _load_files_index() -> list[dict]:
    index_path = DATA_DIR / "files_index.json"
    if index_path.exists():
        return json.loads(index_path.read_text()).get("files", [])
    return []


def _save_files_index(files: list[dict]):
    DATA_DIR.mkdir(parents=True, exist_ok=True)
    index_path = DATA_DIR / "files_index.json"
    index_path.write_text(json.dumps({"files": files}, indent=2, ensure_ascii=False))


def add_internal_doc(title: str, content: str, product: str, tags: list[str] = None) -> dict:
    """Admin: add or update an internal deployment document (text). Same title will be replaced."""
    docs = _load_internal_docs()
    docs = [d for d in docs if d["title"] != title]
    entry = {
        "title": title,
        "content": content,
        "product": product,
        "tags": tags or [],
    }
    docs.append(entry)
    _save_internal_docs(docs)
    return entry


def add_techdoc_file(filename: str, content: bytes, title: str, product: str, description: str = "") -> dict:
    """Admin: upload a datasheet/doc/PDF to techdocs library."""
    FILES_DIR.mkdir(parents=True, exist_ok=True)
    files = _load_files_index()

    # Upsert: remove existing with same title
    old_files = [f for f in files if f.get("title", "").lower() == title.lower()]
    for old in old_files:
        old_path = FILES_DIR / old["stored_name"]
        if old_path.exists():
            old_path.unlink()
    files = [f for f in files if f.get("title", "").lower() != title.lower()]

    file_hash = hashlib.md5(content).hexdigest()[:8]
    safe_name = re.sub(r"[^\w.-]", "-", filename)[:80]
    stored_name = f"{file_hash}-{safe_name}"

    (FILES_DIR / stored_name).write_bytes(content)

    # Extract text for search
    text_content = _extract_text(FILES_DIR / stored_name, filename)

    entry = {
        "title": title,
        "original_name": filename,
        "stored_name": stored_name,
        "product": product,
        "description": description,
        "size_bytes": len(content),
        "uploaded_at": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime()),
        "text_content": text_content[:50000] if text_content else "",
    }
    files.append(entry)
    _save_files_index(files)
    return entry


def _extract_text(file_path: Path, original_name: str) -> str:
    """Try to extract text from uploaded file."""
    ext = Path(original_name).suffix.lower()

    if ext in (".txt", ".md", ".csv"):
        try:
            return file_path.read_text(errors="ignore")
        except Exception:
            return ""

    if ext in (".pptx",):
        try:
            from pptx import Presentation
            prs = Presentation(str(file_path))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                texts.append(text)
            return "\n".join(texts)
        except Exception:
            return ""

    if ext in (".docx",):
        try:
            import zipfile
            soup = BeautifulSoup
            with zipfile.ZipFile(str(file_path)) as z:
                with z.open("word/document.xml") as f:
                    s = BeautifulSoup(f.read(), "html.parser")
                    return s.get_text("\n", strip=True)
        except Exception:
            return ""

    return ""


def list_techdoc_files() -> list[dict]:
    """List all uploaded techdoc files."""
    files = _load_files_index()
    return [{"title": f["title"], "original_name": f["original_name"],
             "product": f["product"], "stored_name": f["stored_name"],
             "description": f.get("description", ""),
             "uploaded_at": f.get("uploaded_at", "")}
            for f in files]


def remove_techdoc_file(stored_name: str) -> bool:
    """Remove an uploaded techdoc file."""
    files = _load_files_index()
    new_files = [f for f in files if f["stored_name"] != stored_name]
    if len(new_files) == len(files):
        return False
    _save_files_index(new_files)
    file_path = FILES_DIR / stored_name
    if file_path.exists():
        file_path.unlink()
    return True


def remove_internal_doc(title: str) -> bool:
    """Admin: remove an internal doc by title."""
    docs = _load_internal_docs()
    new_docs = [d for d in docs if d["title"] != title]
    if len(new_docs) == len(docs):
        return False
    _save_internal_docs(new_docs)
    return True


def _search_official_techdocs(query: str, product: str, client: httpx.Client) -> list[dict]:
    """Search docs.paloaltonetworks.com."""
    results = []
    try:
        resp = client.get(f"{TECHDOCS_BASE}/search", params={"q": query}, follow_redirects=True)
        if resp.status_code == 200:
            soup = BeautifulSoup(resp.text, "html.parser")
            for item in soup.find_all("a", href=True)[:20]:
                href = item["href"]
                text = item.get_text(strip=True)
                if text and "/docs/" in href or "techdocs" in href:
                    full_url = href if href.startswith("http") else f"{TECHDOCS_BASE}{href}"
                    results.append({"title": text, "url": full_url})
    except Exception:
        pass
    return results[:10]


def _search_internal_docs(query: str, product: str) -> list[dict]:
    """Search internal documentation library (text entries)."""
    docs = _load_internal_docs()
    query_lower = query.lower()
    keywords = query_lower.split()

    matches = []
    for doc in docs:
        if product and product != "all" and doc.get("product", "").lower() != product:
            continue
        searchable = f"{doc['title']} {doc.get('content', '')} {' '.join(doc.get('tags', []))}".lower()
        if query_lower in searchable or any(kw in searchable for kw in keywords):
            matches.append(doc)

    return matches


def _search_uploaded_files(query: str, product: str) -> list[dict]:
    """Search uploaded techdoc files."""
    files = _load_files_index()
    query_lower = query.lower()
    keywords = query_lower.split()

    matches = []
    for f in files:
        if product and product != "all" and f.get("product", "").lower() != product:
            continue
        searchable = f"{f.get('title', '')} {f.get('description', '')} {f.get('original_name', '')} {f.get('text_content', '')}".lower()
        if query_lower in searchable or any(kw in searchable for kw in keywords):
            matches.append(f)

    return matches


def _extract_relevant_section(text: str, query: str, context_chars: int = 2000) -> str:
    """Extract the most relevant section from document text."""
    lines = text.split("\n")
    query_keywords = query.lower().split()

    scored_lines = []
    for i, line in enumerate(lines):
        line_lower = line.lower()
        score = sum(1 for kw in query_keywords if kw in line_lower)
        if score > 0:
            scored_lines.append((i, score))

    if not scored_lines:
        return text[:context_chars] + ("..." if len(text) > context_chars else "")

    scored_lines.sort(key=lambda x: x[1], reverse=True)
    best_line_idx = scored_lines[0][0]

    start = max(0, best_line_idx - 5)
    end = min(len(lines), best_line_idx + 20)
    section = "\n".join(lines[start:end])

    if len(section) > context_chars:
        section = section[:context_chars] + "..."

    return section


async def handle(arguments: dict) -> str:
    """Execute the techdocs search skill."""
    query = arguments.get("query", "").strip()
    product = arguments.get("product", "all")

    if not query:
        return "Please provide a deployment or technical question."

    # Search uploaded files first
    file_results = _search_uploaded_files(query, product)

    # Search internal text docs
    internal_results = _search_internal_docs(query, product)

    # Search official TechDocs
    official_results = []
    with httpx.Client(headers=HEADERS, timeout=15) as client:
        official_results = _search_official_techdocs(query, product, client)

    lines = []

    if file_results:
        lines.append("## Uploaded Technical Documents\n")
        for f in file_results[:3]:
            lines.append(f"### {f['title']} ({f['original_name']})")
            lines.append(f"Product: {f.get('product', 'N/A')}")
            if f.get("description"):
                lines.append(f"Description: {f['description']}")
            text = f.get("text_content", "")
            if text:
                relevant = _extract_relevant_section(text, query)
                lines.append(f"\n{relevant}\n")
            else:
                lines.append("\n(Binary file - text extraction unavailable)\n")
            lines.append("---")

    if internal_results:
        lines.append("\n## Internal Deployment Docs\n")
        for doc in internal_results[:5]:
            lines.append(f"### {doc['title']}")
            lines.append(f"Product: {doc.get('product', 'N/A')} | Tags: {', '.join(doc.get('tags', []))}")
            content = doc.get("content", "")
            if len(content) > 500:
                content = content[:500] + "..."
            lines.append(f"\n{content}\n")
            lines.append("---")

    if official_results:
        lines.append("\n## Official TechDocs\n")
        for r in official_results[:8]:
            lines.append(f"- [{r['title']}]({r['url']})")

    if not file_results and not internal_results and not official_results:
        lines.append(
            f"No documentation found for '{query}'.\n\n"
            f"Suggestions:\n"
            f"- Check docs.paloaltonetworks.com directly\n"
            f"- Try broader keywords (e.g., 'GlobalProtect' instead of 'GP portal config')\n"
            f"- Ask admin to upload relevant documents"
        )

    return "\n".join(lines)
