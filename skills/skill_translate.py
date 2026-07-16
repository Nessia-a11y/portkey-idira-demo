"""Skill 7: Document Translation

对 PPTX/DOCX/PDF 文件进行翻译，维持原始排版和字体。
支持中/英/日/韩互译。
"""

import json
import time
import hashlib
import copy
from pathlib import Path
from typing import Optional

DATA_DIR = Path(__file__).parent.parent / "data" / "translations"
OUTPUT_DIR = DATA_DIR / "output"

TOOL_DEFINITION = {
    "type": "function",
    "function": {
        "name": "translate_slide",
        "description": (
            "Translate slide/PPT/document content between languages. "
            "Supports PPTX, DOCX, and PDF files. Maintains original formatting and fonts. "
            "Use this when a user asks to translate a presentation, slide deck, or document. "
            "The user should upload the file through the chat upload feature first."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "source_file": {
                    "type": "string",
                    "description": "The filename of the uploaded file to translate",
                },
                "target_language": {
                    "type": "string",
                    "enum": ["zh", "en", "ja", "ko"],
                    "description": "Target language: 'zh' Chinese, 'en' English, 'ja' Japanese, 'ko' Korean",
                },
                "content_to_translate": {
                    "type": "string",
                    "description": "Direct text content to translate (if no file uploaded)",
                },
            },
            "required": ["target_language"],
        },
    },
}

LANG_NAMES = {
    "zh": "中文",
    "en": "English",
    "ja": "日本語",
    "ko": "한국어",
}


def _ensure_dirs():
    DATA_DIR.mkdir(parents=True, exist_ok=True)
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)


def translate_pptx(file_path: Path, translations: dict[str, str], output_path: Path):
    """Translate PPTX file maintaining formatting. translations = {original: translated}"""
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation(str(file_path))

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        original_text = run.text.strip()
                        if original_text and original_text in translations:
                            run.text = translations[original_text]
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                original_text = run.text.strip()
                                if original_text and original_text in translations:
                                    run.text = translations[original_text]

    prs.save(str(output_path))


def translate_docx(file_path: Path, translations: dict[str, str], output_path: Path):
    """Translate DOCX file maintaining formatting."""
    from docx import Document

    doc = Document(str(file_path))

    for para in doc.paragraphs:
        for run in para.runs:
            original_text = run.text.strip()
            if original_text and original_text in translations:
                run.text = translations[original_text]

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    for run in para.runs:
                        original_text = run.text.strip()
                        if original_text and original_text in translations:
                            run.text = translations[original_text]

    doc.save(str(output_path))


def extract_texts_pptx(file_path: Path) -> list[str]:
    """Extract all translatable text segments from PPTX."""
    from pptx import Presentation

    texts = []
    prs = Presentation(str(file_path))

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        text = run.text.strip()
                        if text and len(text) > 0:
                            texts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                text = run.text.strip()
                                if text and len(text) > 0:
                                    texts.append(text)

    return list(dict.fromkeys(texts))


def extract_texts_docx(file_path: Path) -> list[str]:
    """Extract all translatable text segments from DOCX."""
    from docx import Document

    texts = []
    doc = Document(str(file_path))

    for para in doc.paragraphs:
        for run in para.runs:
            text = run.text.strip()
            if text:
                texts.append(text)

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    for run in para.runs:
                        text = run.text.strip()
                        if text:
                            texts.append(text)

    return list(dict.fromkeys(texts))


def extract_texts_pdf(file_path: Path) -> list[str]:
    """Extract text from PDF (page by page)."""
    try:
        from pypdf import PdfReader
        reader = PdfReader(str(file_path))
        texts = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                for line in text.split("\n"):
                    line = line.strip()
                    if line:
                        texts.append(line)
        return list(dict.fromkeys(texts))
    except Exception:
        return []


def build_translated_pdf(file_path: Path, translations: dict[str, str], output_path: Path):
    """Create translated PDF. Best-effort: extracts text and creates new PDF with translations."""
    try:
        from pypdf import PdfReader
        from reportlab.lib.pagesizes import A4
        from reportlab.pdfgen import canvas
        from reportlab.lib.units import inch
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont

        # Try to register a CJK font
        try:
            pdfmetrics.registerFont(TTFont('NotoSansCJK', '/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc'))
            font_name = 'NotoSansCJK'
        except Exception:
            try:
                pdfmetrics.registerFont(TTFont('NotoSans', '/usr/share/fonts/truetype/noto/NotoSans-Regular.ttf'))
                font_name = 'NotoSans'
            except Exception:
                font_name = 'Helvetica'

        reader = PdfReader(str(file_path))
        c = canvas.Canvas(str(output_path), pagesize=A4)
        width, height = A4

        for page in reader.pages:
            text = page.extract_text()
            if not text:
                continue

            c.setFont(font_name, 10)
            y = height - 50
            for line in text.split("\n"):
                line = line.strip()
                if not line:
                    y -= 14
                    continue

                translated_line = translations.get(line, line)
                if y < 50:
                    c.showPage()
                    c.setFont(font_name, 10)
                    y = height - 50

                c.drawString(40, y, translated_line)
                y -= 14

            c.showPage()

        c.save()
    except ImportError:
        # Fallback: write translated text as plain text
        lines = []
        try:
            from pypdf import PdfReader
            reader = PdfReader(str(file_path))
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    for line in text.split("\n"):
                        line = line.strip()
                        if line:
                            translated = translations.get(line, line)
                            lines.append(translated)
                    lines.append("---")
        except Exception:
            lines = ["PDF translation failed - missing dependencies"]

        output_txt = output_path.with_suffix(".txt")
        output_txt.write_text("\n".join(lines), encoding="utf-8")


async def handle(arguments: dict) -> str:
    """Execute the translation skill — returns instructions for the agent."""
    source_file = arguments.get("source_file", "").strip()
    target_language = arguments.get("target_language", "zh").strip()
    content_to_translate = arguments.get("content_to_translate", "").strip()

    target_lang_name = LANG_NAMES.get(target_language, target_language)

    if content_to_translate:
        return (
            f"TRANSLATION_REQUEST:\n"
            f"Target language: {target_lang_name}\n"
            f"Content ({len(content_to_translate)} chars):\n\n"
            f"{content_to_translate}\n\n"
            f"---\n"
            f"Translate the above to {target_lang_name}. "
            f"Keep original structure. Technical terms stay in English with {target_lang_name} notes in parentheses."
        )

    if source_file:
        # Check uploaded files
        upload_dir = Path(__file__).parent.parent / "data" / "uploads"
        file_path = upload_dir / source_file
        if not file_path.exists():
            # Search by partial name
            if upload_dir.exists():
                for f in upload_dir.iterdir():
                    if source_file.lower() in f.name.lower():
                        file_path = f
                        break

        if file_path.exists():
            ext = file_path.suffix.lower()
            texts = []

            if ext == ".pptx":
                texts = extract_texts_pptx(file_path)
            elif ext == ".docx":
                texts = extract_texts_docx(file_path)
            elif ext == ".pdf":
                texts = extract_texts_pdf(file_path)
            else:
                return f"Unsupported file format: {ext}. Supported: .pptx, .docx, .pdf"

            if not texts:
                return f"Could not extract text from {source_file}. The file may be empty or in an unsupported format."

            # Return text for agent to translate
            text_block = "\n---SEGMENT---\n".join(texts[:200])
            return (
                f"FILE_TRANSLATION_REQUEST:\n"
                f"File: {file_path.name}\n"
                f"Format: {ext}\n"
                f"Target language: {target_lang_name}\n"
                f"Segments to translate: {len(texts[:200])}\n\n"
                f"IMPORTANT: Translate each segment below to {target_lang_name}. "
                f"Return the result as a JSON object with original text as keys and translations as values. "
                f"Keep formatting markers (bullets, numbering) intact. "
                f"Technical terms (product names, protocols) stay in English.\n\n"
                f"SEGMENTS:\n{text_block}"
            )

        return (
            f"File '{source_file}' not found. "
            f"Please ask the user to upload the file first using the upload button in the chat interface."
        )

    return (
        f"No file or content specified for translation.\n"
        f"Please ask the user to upload a file (.pptx, .docx, or .pdf) "
        f"or provide text content directly."
    )
